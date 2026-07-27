#!/usr/bin/env python3
"""Extract text, notes, and images from a .pptx file into a simple folder."""
from __future__ import annotations

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError as exc:
    raise SystemExit("Missing dependency: install python-pptx with `pip install python-pptx`") from exc


def shape_text(shape) -> str:
    if not hasattr(shape, "text"):
        return ""
    return (shape.text or "").strip()


def extract(input_pptx: Path, output_dir: Path) -> None:
    prs = Presentation(str(input_pptx))
    assets_dir = output_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    slides = []

    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        images = []
        for s_idx, shape in enumerate(slide.shapes, start=1):
            text = shape_text(shape)
            if text:
                texts.append(text)
            if getattr(shape, "shape_type", None) == 13:  # picture
                image = shape.image
                ext = image.ext or "png"
                filename = f"slide-{idx:02d}-image-{s_idx:02d}.{ext}"
                path = assets_dir / filename
                path.write_bytes(image.blob)
                images.append(str(path.relative_to(output_dir)))
        notes = ""
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes = ""
        slides.append({"slide": idx, "texts": texts, "images": images, "notes": notes})

    output_dir.mkdir(parents=True, exist_ok=True)
    (output_dir / "slides.json").write_text(json.dumps(slides, ensure_ascii=False, indent=2), encoding="utf-8")
    md = [f"# Extracted PPTX: {input_pptx.name}", ""]
    for slide in slides:
        md.append(f"## Slide {slide['slide']}")
        for text in slide["texts"]:
            md.append(f"- {text.replace(chr(10), ' / ')}")
        if slide["images"]:
            md.append(f"Images: {', '.join(slide['images'])}")
        if slide["notes"]:
            md.append(f"Notes: {slide['notes']}")
        md.append("")
    (output_dir / "slides.md").write_text("\n".join(md), encoding="utf-8")


def main() -> None:
    if len(sys.argv) != 3:
        raise SystemExit("Usage: python scripts/extract-pptx.py <input.pptx> <output_dir>")
    input_pptx = Path(sys.argv[1]).expanduser().resolve()
    output_dir = Path(sys.argv[2]).expanduser().resolve()
    if not input_pptx.exists():
        raise SystemExit(f"Input not found: {input_pptx}")
    extract(input_pptx, output_dir)
    print(f"Extracted to {output_dir}")


if __name__ == "__main__":
    main()
